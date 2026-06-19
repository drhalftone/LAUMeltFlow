"""Generate ECE Symposium 2026 poster as a .pptx file."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = Path(__file__).parent

# ------------------------------------------------------------------
# Poster geometry: 48" wide x 36" tall (standard ECE symposium)
# ------------------------------------------------------------------
POSTER_W = Inches(48)
POSTER_H = Inches(36)

MARGIN = Inches(0.6)
HEADER_H = Inches(4.2)
FOOTER_H = Inches(1.0)
GUTTER = Inches(0.5)

# Column geometry
COL_Y = MARGIN + HEADER_H + Inches(0.3)
COL_H = POSTER_H - COL_Y - FOOTER_H - Inches(0.3)
COL_W = (POSTER_W - 2 * MARGIN - 2 * GUTTER) / 3

# Colors (University of Kentucky blue + neutral)
UK_BLUE = RGBColor(0x00, 0x33, 0xA0)
UK_BLUE_DARK = RGBColor(0x00, 0x26, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
ACCENT = RGBColor(0xC8, 0x9B, 0x2A)  # UK accent gold


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             bg=None, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    if bg is not None:
        set_fill(box, bg)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=16, color=BLACK,
                space_after=4, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, color, line=False):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    if line:
        r.line.color.rgb = UK_BLUE_DARK
        r.line.width = Pt(1.0)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def add_block_header(slide, x, y, w, title, *, height=Inches(0.7)):
    add_rect(slide, x, y, w, height, UK_BLUE)
    add_text(slide, x, y, w, height, title,
             size=28, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font="Calibri")
    return y + height


def add_image_fit(slide, path, x, y, w, h):
    """Add image scaled to fit within box while preserving aspect ratio."""
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    aspect = iw / ih
    box_aspect = w / h
    if aspect > box_aspect:
        nw = w
        nh = int(w / aspect)
    else:
        nh = h
        nw = int(h * aspect)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    slide.shapes.add_picture(str(path), nx, ny, width=nw, height=nh)


def add_column_background(slide, x, y, w, h):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = LIGHT
    r.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    r.line.width = Pt(0.5)
    r.shadow.inherit = False
    return r


# ------------------------------------------------------------------
# Build presentation
# ------------------------------------------------------------------
prs = Presentation()
prs.slide_width = POSTER_W
prs.slide_height = POSTER_H

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# Background
bg = add_rect(slide, 0, 0, POSTER_W, POSTER_H, WHITE)

# ------------------------------------------------------------------
# Header bar
# ------------------------------------------------------------------
add_rect(slide, 0, 0, POSTER_W, HEADER_H, UK_BLUE)
# Accent stripe under header
add_rect(slide, 0, HEADER_H, POSTER_W, Inches(0.18), ACCENT)

add_text(
    slide, MARGIN, Inches(0.45), POSTER_W - 2 * MARGIN, Inches(2.2),
    "Message-Passing Graph Neural Network as a Surrogate for\n"
    "Variable-Length 1D Finite Element Bead Chain Simulations",
    size=66, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
)
add_text(
    slide, MARGIN, Inches(2.75), POSTER_W - 2 * MARGIN, Inches(0.7),
    "Grey L. Goodwin    |    Daniel L. Lau, PhD.",
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
)
add_text(
    slide, MARGIN, Inches(3.35), POSTER_W - 2 * MARGIN, Inches(0.8),
    "Department of Electrical and Computer Engineering  ·  University of Kentucky  ·  "
    "grey.goodwin@uky.edu  ·  dllau@uky.edu",
    size=24, color=WHITE, align=PP_ALIGN.CENTER,
)

# ------------------------------------------------------------------
# Column backgrounds
# ------------------------------------------------------------------
col_x = [
    MARGIN,
    MARGIN + COL_W + GUTTER,
    MARGIN + 2 * (COL_W + GUTTER),
]
for cx in col_x:
    add_column_background(slide, cx, COL_Y, COL_W, COL_H)

# ------------------------------------------------------------------
# Column 1: Abstract · Problem Setup · MPGNN Architecture
# ------------------------------------------------------------------
x = col_x[0]
w_inner = COL_W - Inches(0.4)
x_inner = x + Inches(0.2)
y = COL_Y + Inches(0.2)

# --- Abstract ---
y = add_block_header(slide, x_inner, y, w_inner, "Abstract")
abstract = (
    "We train a message-passing graph neural network (BeadMPGNN) as a per-timestep "
    "surrogate for a 1D finite element bead chain. A key finding: even a well-designed "
    "model fails without the right training data. Sampling neighbor positions in polar "
    "coordinates near the rest length was necessary for stable rollouts, while uniform "
    "Cartesian sampling produced models that diverged.\n\n"
    "We restructure an initial fixed-size MLP as a true MPGNN with separate node and "
    "edge encoders, ghost masking for endpoints, and weight sharing across beads. One "
    "trained model runs on chains of any length. Trained on 1M volume-sampled per-bead "
    "scenarios spanning 8–32 beads and random tapers, it generalizes to chain lengths, "
    "orientations, and velocities unseen during training."
)
box = add_text(slide, x_inner, y, w_inner, Inches(6.5), abstract,
               size=17, color=BLACK)
y = y + Inches(6.6)

# --- Problem Setup ---
y = add_block_header(slide, x_inner, y, w_inner, "Problem Setup")
problem = (
    "• Chain of N mass nodes connected by N−1 spring rod elements.\n"
    "• Tapered mass distribution from heavy anchor to light tip (ratio r = m₀/m_{N−1}).\n"
    "• Forces per bead: gravity, Hooke's-law spring forces from each connected neighbor, "
    "viscous rod damping, and a small global velocity drag.\n"
    "• Reference solver uses symplectic Euler integration (Δt = 1e−4 s).\n"
    "• Physics is strictly local: force on bead i depends only on i and its immediate "
    "neighbors i−1, i+1."
)
add_text(slide, x_inner, y, w_inner, Inches(4.2), problem, size=17, color=BLACK)
y = y + Inches(4.3)

# --- MPGNN Architecture ---
y = add_block_header(slide, x_inner, y, w_inner, "MPGNN Architecture")
arch_text = (
    "Five learned components, all weight-shared across beads:\n"
    "• Node encoder → 64-d hidden from 4-d node features\n"
    "• Edge encoder → 64-d hidden from 6-d rel. edge features\n"
    "• Message MLP → per-edge 64-d message\n"
    "• Node update → [h_self ∥ m_left ∥ m_right] → 64-d\n"
    "• Output head → 4-d Δstate (Δx, Δy, Δvx, Δvy)\n"
    "Ghost masking zeros messages from absent endpoint neighbors. 80,324 params; the "
    "same model runs on any chain length."
)
add_text(slide, x_inner, y, w_inner, Inches(4.4), arch_text, size=17, color=BLACK)
y = y + Inches(4.5)

# Architecture figure
arch_img = HERE / "architecture_mpgnn_only.png"
img_h = (COL_Y + COL_H) - y - Inches(0.2)
add_image_fit(slide, arch_img, x_inner, y, w_inner, img_h)

# ------------------------------------------------------------------
# Column 2: Pipeline · Why Volume Sampling · Acknowledgements
# ------------------------------------------------------------------
x = col_x[1]
x_inner = x + Inches(0.2)
y = COL_Y + Inches(0.2)

# --- Pipeline figure ---
y = add_block_header(slide, x_inner, y, w_inner, "Volume-Sampled Training Pipeline")
pipeline_img = HERE / "pipeline_visualization.png"
pipeline_h = Inches(9.5)
add_image_fit(slide, pipeline_img, x_inner, y, w_inner, pipeline_h)
y = y + pipeline_h + Inches(0.05)
add_text(slide, x_inner, y, w_inner, Inches(1.3),
         "Each sample is an independent random bead-level scenario: random N ∈ [8, 32], "
         "taper ratio r ∈ [1, 10], polar-sampled neighbors near rest length. Physics "
         "solved analytically; no trajectories recorded. Dataset: 1,000,000 samples.",
         size=15, color=GRAY)
y = y + Inches(1.4)

# --- Why Volume Sampling ---
y = add_block_header(slide, x_inner, y, w_inner, "Why Volume Sampling")
why_text = (
    "Trajectory-recorded data traces a thin path through state space. During rollout, "
    "small prediction errors push the chain off this path into states the model has "
    "never seen — errors compound and the simulation diverges.\n\n"
    "Volume sampling fills the reachable state space uniformly. Every training sample "
    "has a physically realistic rod length (±5% of ℓ₀), so the model learns force "
    "magnitudes matched to what it sees at inference.\n\n"
    "• Cartesian uniform sampling → stable loss, diverging rollouts\n"
    "• Polar near rest length → stable, accurate rollouts\n"
    "Sampling strategy matters as much as architecture."
)
add_text(slide, x_inner, y, w_inner, Inches(6.0), why_text, size=17, color=BLACK)
y = y + Inches(6.1)

# --- Acknowledgements ---
y = add_block_header(slide, x_inner, y, w_inner, "Acknowledgements")
ack_text = (
    "This work was supported in part by the National Science Foundation under Grants "
    "CCF 2230161 and 2230162, and in part by AFOSR under Grant FA9550-22-1-0362."
)
add_text(slide, x_inner, y, w_inner, Inches(2.0), ack_text, size=16, color=BLACK)

# ------------------------------------------------------------------
# Column 3: Rollout · Training loss · Generalization · Conclusions
# ------------------------------------------------------------------
x = col_x[2]
x_inner = x + Inches(0.2)
y = COL_Y + Inches(0.2)

# --- Rollout snapshot ---
y = add_block_header(slide, x_inner, y, w_inner, "Side-by-Side Rollout")
roll_img = HERE / "frame_falling.png"
roll_h = Inches(5.2)
add_image_fit(slide, roll_img, x_inner, y, w_inner, roll_h)
y = y + roll_h + Inches(0.05)
add_text(slide, x_inner, y, w_inner, Inches(1.1),
         "Reference solver (black, left) vs. BeadMPGNN (red, right) on a 16-bead chain "
         "released from horizontal. The GNN tracks the reference through fall, swing, "
         "and tip whip.",
         size=15, color=GRAY)
y = y + Inches(1.2)

# --- Training loss ---
y = add_block_header(slide, x_inner, y, w_inner, "Training Convergence")
loss_img = HERE / "mpgnn_loss_only.png"
loss_h = Inches(4.0)
add_image_fit(slide, loss_img, x_inner, y, w_inner, loss_h)
y = y + loss_h + Inches(0.05)
add_text(slide, x_inner, y, w_inner, Inches(0.9),
         "Train/validation MSE over 200 epochs. AdamW, lr = 1e−3 with cosine anneal, "
         "batch 512, 10% validation split, gradient clipping at norm 1.0.",
         size=15, color=GRAY)
y = y + Inches(1.0)

# --- Generalization ---
y = add_block_header(slide, x_inner, y, w_inner, "Generalization Across N")
gen_img = HERE / "mpgnn_pos_error_only.png"
gen_h = Inches(4.0)
add_image_fit(slide, gen_img, x_inner, y, w_inner, gen_h)
y = y + gen_h + Inches(0.05)
add_text(slide, x_inner, y, w_inner, Inches(1.0),
         "Mean position error vs. reference for N ∈ {8, 10, 12, 16, 20, 24, 32}. Short "
         "to medium chains track closely; longer chains show expected drift from "
         "compounding per-step errors.",
         size=15, color=GRAY)
y = y + Inches(1.1)

# --- Conclusions ---
y = add_block_header(slide, x_inner, y, w_inner, "Conclusions")
concl = (
    "• Sampling strategy is as important as architecture. Polar sampling near rest "
    "length is required; Cartesian uniform sampling fails.\n"
    "• MPGNN with weight sharing + ghost masking enables one model across chain lengths.\n"
    "• Volume sampling covers reachable state space — no trajectory recording needed.\n"
    "• Model generalizes to unseen N, orientations, and velocities within sampling bounds.\n"
    "• Future: rollout training to reduce drift; K > 1 message passing for longer-range "
    "effects."
)
add_text(slide, x_inner, y, w_inner,
         (COL_Y + COL_H) - y - Inches(0.2), concl, size=16, color=BLACK)

# ------------------------------------------------------------------
# Footer
# ------------------------------------------------------------------
fy = POSTER_H - FOOTER_H
add_rect(slide, 0, fy, POSTER_W, FOOTER_H, UK_BLUE_DARK)
add_text(slide, MARGIN, fy, POSTER_W - 2 * MARGIN, FOOTER_H,
         "ECE Symposium 2026  ·  Department of Electrical and Computer Engineering  ·  "
         "University of Kentucky",
         size=22, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# ------------------------------------------------------------------
out = HERE / "ECE_Symposium_Poster_2026.pptx"
prs.save(out)
print(f"Wrote {out}")
